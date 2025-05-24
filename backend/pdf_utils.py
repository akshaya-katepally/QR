# import fitz  # PyMuPDF

# def extract_texts(files, raw_text):
#     text = raw_text or ''
#     for file in files:
#         if file and file.filename.endswith('.pdf'):
#             doc = fitz.open(stream=file.read(), filetype="pdf")
#             for page in doc:
#                 text += page.get_text()
#     return text.strip()

import fitz  # PyMuPDF
import docx
import pptx
import pytesseract
from PIL import Image
import io
import pandas as pd

def extract_texts(files, raw_text):
    text = raw_text.strip() + "\n" if raw_text else ""

    for file in files:
        filename = file.filename.lower()
        content = file.read()

        if filename.endswith(".pdf"):
            try:
                doc = fitz.open(stream=content, filetype="pdf")
                for page in doc:
                    text += page.get_text() + "\n"
            except Exception as e:
                text += f"\n[Error reading PDF: {str(e)}]\n"

        elif filename.endswith(".docx"):
            try:
                doc = docx.Document(io.BytesIO(content))
                for para in doc.paragraphs:
                    text += para.text + "\n"
            except Exception as e:
                text += f"\n[Error reading DOCX: {str(e)}]\n"

        elif filename.endswith(".pptx"):
            try:
                presentation = pptx.Presentation(io.BytesIO(content))
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            except Exception as e:
                text += f"\n[Error reading PPTX: {str(e)}]\n"

        elif filename.endswith(".xlsx"):
            try:
                excel_file = pd.ExcelFile(io.BytesIO(content))
                for sheet_name in excel_file.sheet_names:
                    df = excel_file.parse(sheet_name)
                    text += f"--- Sheet: {sheet_name} ---\n"
                    text += df.astype(str).to_string(index=False) + "\n"
            except Exception as e:
                text += f"\n[Error reading XLSX: {str(e)}]\n"

        elif filename.endswith(".txt"):
            try:
                text += content.decode("utf-8", errors="ignore") + "\n"
            except Exception as e:
                text += f"\n[Error reading TXT: {str(e)}]\n"

        elif filename.endswith(('.jpg', '.jpeg', '.png', '.bmp', '.tiff')):
            try:
                img = Image.open(io.BytesIO(content))
                text += pytesseract.image_to_string(img) + "\n"
            except Exception as e:
                text += f"\n[Error reading image: {str(e)}]\n"

        else:
            text += f"\n[Unsupported file type: {filename}]\n"

    return text.strip()
